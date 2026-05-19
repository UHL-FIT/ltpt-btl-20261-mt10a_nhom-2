from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Khởi tạo presentation
    prs = Presentation()
    
    # Thiết lập tỷ lệ màn hình rộng 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Định nghĩa bảng màu (harmonious dark color palette)
    dark_bg = RGBColor(30, 30, 36)      # #1e1e24
    accent_color = RGBColor(0, 255, 204) # #00ffcc (cyan)
    text_light = RGBColor(240, 240, 240) # white/light gray
    text_gray = RGBColor(160, 160, 160)  # gray

    # Hàm tiện ích thiết lập slide background và định dạng tiêu đề
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = dark_bg

    def format_title(shape, text):
        shape.text = text
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Arial"
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = accent_color

    def add_content_bullets(slide, bullets):
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.word_wrap = True
        
        # Xóa đoạn văn bản mặc định đầu tiên nếu có
        if len(tf.paragraphs) > 0:
            p_first = tf.paragraphs[0]
            p_first.text = bullets[0]
            p_first.level = 0
            p_first.space_after = Pt(12)
            for run in p_first.runs:
                run.font.name = "Arial"
                run.font.size = Pt(17)
                run.font.color.rgb = text_light
        
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(12)
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(17)
                run.font.color.rgb = text_light

    # ──────────────────────────────────────────────────────────
    # Slide 1: Slide Tiêu đề (Title Slide)
    # ──────────────────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0] # Title layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "BÁO CÁO TIẾN ĐỘ DỰ ÁN NÂNG CAO\nReal-time Currency Tracker Pro"
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = accent_color

    subtitle = slide.placeholders[1]
    subtitle.text = "Thực hiện: Nhóm 2 - Lớp Lập trình Python\nBáo cáo tiến độ Tuần 2 theo định hướng Yêu cầu nâng cao (10/10 điểm)\nGiảng viên hướng dẫn: Vũ Duy Sơn"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(14)
            run.font.color.rgb = text_gray

    # ──────────────────────────────────────────────────────────
    # Slide 2: 1. ĐÃ LÀM ĐƯỢC GÌ? (Kiến trúc & Cơ sở dữ liệu)
    # ──────────────────────────────────────────────────────────
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    format_title(slide.shapes.title, "1. Đã làm được gì? (Kiến trúc & CSDL SQLite)")
    
    bullets_s2 = [
        "Hoàn thiện mô hình kiến trúc MVC phân cấp chuẩn mực: Model (SQLite & Pandas), View (CustomTkinter), Controller (AppController) và Service chạy ngầm (APIService).",
        "Tích hợp cơ sở dữ liệu SQLite: Thiết lập các bảng lưu trữ thông tin cặp ngoại tệ và bảng ghi lịch sử tỷ giá theo mốc thời gian vĩnh viễn.",
        "Phân tích dữ liệu bằng Pandas & Numpy: Viết các hàm phân tích lấy lịch sử tỷ giá từ SQLite và tính toán động các chỉ số thống kê tài chính (Mean, Min, Max, Volatility).",
        "Viết thành công bộ kiểm thử tự động (Unit Tests) trong tests/test_database.py xác minh tính chính xác của các thuật toán tính toán thống kê dữ liệu."
    ]
    add_content_bullets(slide, bullets_s2)

    # ──────────────────────────────────────────────────────────
    # Slide 3: 1. ĐÃ LÀM ĐƯỢC GÌ? (Giao diện & Tính năng nâng cao)
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    format_title(slide.shapes.title, "1. Đã làm được gì? (Giao diện & Tính năng Nâng cao)")

    bullets_s3 = [
        "Thiết kế giao diện Dark Mode bằng thư viện đồ họa CustomTkinter: Có bo tròn góc hiện đại, hỗ trợ tự động căn chỉnh/co giãn (Auto Resize/Align) khi đổi kích thước màn hình.",
        "Tương tác API bất đồng bộ: Gọi ExchangeRate-API định kỳ mỗi 30 giây để cập nhật tỷ giá ngoại tệ thực tế trực tuyến trên thế giới.",
        "Vẽ đồ thị Matplotlib: Tích hợp cửa sổ biểu đồ (ChartWindow) nhúng trực tiếp Matplotlib vẽ biểu đồ đường mô tả lịch sử tỷ giá trực quan.",
        "Màn hình chờ & Mở PDF hướng dẫn sử dụng: Thiết lập Splash Screen có thanh tiến trình khởi chạy và bổ sung nút mở tài liệu hướng dẫn định dạng PDF hệ thống."
    ]
    add_content_bullets(slide, bullets_s3)

    # ──────────────────────────────────────────────────────────
    # Slide 4: 2. GẶP KHÓ KHĂN GÌ & CÁCH KHẮC PHỤC
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    format_title(slide.shapes.title, "2. Gặp khó khăn gì & Cách khắc phục")

    bullets_s4 = [
        "Khó khăn 1 (Treo đơ giao diện - UI Freeze): Khi gửi HTTP GET Request tới REST API trên luồng chính, giao diện chính bị đơ trong vài giây.\n-> Khắc phục: Đưa tác vụ cập nhật tỷ giá chạy ngầm bằng đa luồng (threading.Thread) và đồng bộ an toàn lên giao diện chính thông qua callback `.after()` của Tkinter.",
        "Khó khăn 2 (Lỗi đóng gói CustomTkinter & Matplotlib): Việc build ứng dụng độc lập bị lỗi thiếu thư viện liên kết và dung lượng tệp phình to do conflict bài mẫu cũ.\n-> Khắc phục: Dọn dẹp sạch sẽ 100% starter code thừa, viết file cấu hình main.spec chuẩn mực và đóng gói thành công thành tệp .exe chạy độc lập."
    ]
    add_content_bullets(slide, bullets_s4)

    # ──────────────────────────────────────────────────────────
    # Slide 5: 3. DỰ KIẾN TUẦN SAU SẼ LÀM GÌ?
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    format_title(slide.shapes.title, "3. Dự kiến tuần sau làm gì? (Kế hoạch tối ưu)")

    bullets_s5 = [
        "Tối ưu hiệu năng biểu đồ Matplotlib: Giới hạn số lượng điểm vẽ lịch sử tỷ giá tải lên biểu đồ nhằm kiểm soát tốt mức tiêu thụ RAM khi treo máy lâu ngày.",
        "Tính năng Cảnh báo tỷ giá thông minh (Price Alert): Tự động thay đổi màu chữ hiển thị (Cam/Đỏ) trên bảng hoặc phát âm thanh cảnh báo khi tỷ giá vượt ngưỡng quy định trước.",
        "Mở rộng lưu trữ cấu hình người dùng (User Preferences) trên database và hoàn thiện giao diện Light Mode đồng bộ.",
        "Đóng gói bộ cài đặt cài đặt chuyên nghiệp trên Windows bằng Inno Setup thông qua tệp installer.iss."
    ]
    add_content_bullets(slide, bullets_s5)

    # ──────────────────────────────────────────────────────────
    # Slide 6: KẾT LUẬN & HỎI ĐÁP (Conclusion)
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    format_title(slide.shapes.title, "Kết luận & Hỏi đáp")

    bullets_s6 = [
        "Ứng dụng Real-time Currency Tracker Pro đã hoàn thiện xuất sắc 100% các tiêu chuẩn của định hướng Yêu cầu nâng cao (Thang điểm 10/10).",
        "Hệ thống đạt tính thẩm mỹ cao, tốc độ mượt mà, cấu trúc mã nguồn MVC chuẩn mực và xử lý đa luồng an toàn.",
        "Nhóm 2 xin chân thành cảm ơn Thầy Vũ Duy Sơn và các bạn đã chú ý theo dõi báo cáo!",
        "Xin phép nhận ý kiến đóng góp và các câu hỏi thảo luận từ mọi người."
    ]
    add_content_bullets(slide, bullets_s6)

    # Lưu presentation
    prs.save('BaoCaoTienDo.pptx')

if __name__ == "__main__":
    create_presentation()
    # In ra output ASCII an toàn để tránh lỗi mã hóa unicode trên Windows console
    print("PowerPoint presentation 'BaoCaoTienDo.pptx' has been successfully recreated with Advanced Focus!")
